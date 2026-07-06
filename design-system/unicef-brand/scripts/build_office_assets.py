#!/usr/bin/env python3
"""
Build UNICEF Brand 4.0 Office assets: Theme (.thmx), PowerPoint template (.potx),
and Word template (.dotx).

All values are drawn from UNICEF Brand Book 4.0 (Colour p.49-56, Fonts p.57-61).
The type scale and document spacing are a reasonable authoring extension and are
NOT brand-mandated; see design-system.md for which is which.

Run:  python3 build_office_assets.py
Outputs land in ../office/ next to this script.
"""
import copy
import os
import shutil
import zipfile
from lxml import etree

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

import docx
from docx import Document
from docx.shared import Pt as DPt, RGBColor as DRGB, Inches as DInches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn, nsmap

HERE = os.path.dirname(os.path.abspath(__file__))
OUT = os.path.normpath(os.path.join(HERE, "..", "office"))
os.makedirs(OUT, exist_ok=True)

A = "http://schemas.openxmlformats.org/drawingml/2006/main"

# ---------------------------------------------------------------------------
# Brand tokens (verified from Brand Book 4.0)
# ---------------------------------------------------------------------------
UNICEF_BLUE = "00AEEF"   # Primary — Process Cyan, RGB 0/174/239. Dominant colour.
WHITE = "FFFFFF"
BLACK = "000000"
ACC_BLUE = "0047BB"      # Pantone 2728
ACC_GREEN = "004C45"     # Pantone 3302C
ACC_ORANGE_D = "FF7100"  # Pantone 2728 (dark orange in book)
ACC_ORANGE = "FF8200"    # Pantone 151C
ACC_YELLOW = "FFB500"    # Pantone 7549C
ACC_LTCYAN = "9ADBE8"    # Pantone 304C
EMERGENCY_RED = "E2231A" # Pantone 485

MAJOR_FONT = "Noto Sans"
MINOR_FONT = "Noto Sans"
THEME_NAME = "UNICEF Brand 4.0"


def srgb(parent_tag, hexval):
    el = etree.SubElement(parent_tag, f"{{{A}}}srgbClr")
    el.set("val", hexval)
    return el


def build_clr_scheme():
    ns = f"{{{A}}}"
    clr = etree.Element(f"{ns}clrScheme")
    clr.set("name", THEME_NAME)

    def sys(tag, val, last):
        wrap = etree.SubElement(clr, f"{ns}{tag}")
        s = etree.SubElement(wrap, f"{ns}sysClr")
        s.set("val", val)
        s.set("lastClr", last)

    def rgb(tag, hexval):
        wrap = etree.SubElement(clr, f"{ns}{tag}")
        s = etree.SubElement(wrap, f"{ns}srgbClr")
        s.set("val", hexval)

    sys("dk1", "windowText", BLACK)      # text/background dark 1  -> Black
    sys("lt1", "window", WHITE)          # text/background light 1 -> White
    rgb("dk2", ACC_GREEN)                # dark 2
    rgb("lt2", ACC_LTCYAN)               # light 2
    rgb("accent1", UNICEF_BLUE)          # DOMINANT brand colour
    rgb("accent2", ACC_BLUE)
    rgb("accent3", ACC_GREEN)
    rgb("accent4", ACC_ORANGE)
    rgb("accent5", ACC_YELLOW)
    rgb("accent6", ACC_LTCYAN)
    rgb("hlink", ACC_BLUE)
    rgb("folHlink", ACC_GREEN)
    return clr


def build_font_scheme():
    ns = f"{{{A}}}"
    fs = etree.Element(f"{ns}fontScheme")
    fs.set("name", THEME_NAME)
    for major in ("majorFont", "minorFont"):
        mf = etree.SubElement(fs, f"{ns}{major}")
        latin = etree.SubElement(mf, f"{ns}latin")
        latin.set("typeface", MAJOR_FONT if major == "majorFont" else MINOR_FONT)
        etree.SubElement(mf, f"{ns}ea").set("typeface", "")
        etree.SubElement(mf, f"{ns}cs").set("typeface", "")
    return fs


def make_theme_xml():
    """Start from python-pptx's default theme (valid fmtScheme) and swap in the
    UNICEF colour + font schemes. Returns bytes of a complete theme1.xml."""
    prs = Presentation()
    theme_part = prs.slide_masters[0].part.part_related_by(
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme"
    )
    root = etree.fromstring(theme_part.blob)
    root.set("name", THEME_NAME)
    ns = f"{{{A}}}"
    theme_elements = root.find(f"{ns}themeElements")

    old_clr = theme_elements.find(f"{ns}clrScheme")
    theme_elements.replace(old_clr, build_clr_scheme())
    old_font = theme_elements.find(f"{ns}fontScheme")
    theme_elements.replace(old_font, build_font_scheme())

    return etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)


THEME_XML = make_theme_xml()


# ---------------------------------------------------------------------------
# 1. Standalone Office Theme  (.thmx)
# ---------------------------------------------------------------------------
def build_thmx(path):
    content_types = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">'
        '<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>'
        '<Default Extension="xml" ContentType="application/xml"/>'
        '<Override PartName="/theme/theme/themeManager.xml" '
        'ContentType="application/vnd.openxmlformats-officedocument.themeManager+xml"/>'
        '<Override PartName="/theme/theme/theme1.xml" '
        'ContentType="application/vnd.openxmlformats-officedocument.theme+xml"/>'
        '</Types>'
    )
    root_rels = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
        '<Relationship Id="rId1" '
        'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/themeManager" '
        'Target="theme/theme/themeManager.xml"/>'
        '</Relationships>'
    )
    theme_mgr = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<p:presentation xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
        'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" '
        'xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"/>'
    )
    mgr_rels = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
        '<Relationship Id="rId1" '
        'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme" '
        'Target="theme1.xml"/>'
        '</Relationships>'
    )
    if os.path.exists(path):
        os.remove(path)
    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as z:
        z.writestr("[Content_Types].xml", content_types)
        z.writestr("_rels/.rels", root_rels)
        z.writestr("theme/theme/themeManager.xml", theme_mgr)
        z.writestr("theme/_rels/themeManager.xml.rels", mgr_rels)
        z.writestr("theme/theme/theme1.xml", THEME_XML)
    return path


# ---------------------------------------------------------------------------
# Helper: overwrite the theme part inside an OPC package (pptx/docx) on disk
# ---------------------------------------------------------------------------
def inject_theme(package_path):
    tmp = package_path + ".tmp"
    with zipfile.ZipFile(package_path, "r") as zin:
        names = zin.namelist()
        theme_members = [n for n in names if n.startswith(("ppt/theme/", "word/theme/"))
                         and n.endswith(".xml")]
        with zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                data = zin.read(item.filename)
                if item.filename in theme_members:
                    data = THEME_XML
                zout.writestr(item, data)
    os.replace(tmp, package_path)


# ---------------------------------------------------------------------------
# Convert a saved .pptx/.docx into a template by rewriting the main content type
# ---------------------------------------------------------------------------
def to_template(package_path, template_path, kind):
    mapping = {
        "pptx": (
            "application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml",
            "application/vnd.openxmlformats-officedocument.presentationml.template.main+xml",
        ),
        "docx": (
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml",
            "application/vnd.openxmlformats-officedocument.wordprocessingml.template.main+xml",
        ),
    }
    src_ct, dst_ct = mapping[kind]
    tmp = template_path + ".tmp"
    with zipfile.ZipFile(package_path, "r") as zin:
        with zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                data = zin.read(item.filename)
                if item.filename == "[Content_Types].xml":
                    data = data.replace(src_ct.encode(), dst_ct.encode())
                zout.writestr(item, data)
    os.replace(tmp, template_path)


# ---------------------------------------------------------------------------
# 2. PowerPoint template (.potx)
# ---------------------------------------------------------------------------
def build_potx(path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)   # 16:9
    prs.slide_height = Inches(7.5)

    blue = RGBColor.from_string(UNICEF_BLUE)
    white = RGBColor.from_string(WHITE)
    black = RGBColor.from_string(BLACK)

    def add_footer(slide, dark_bg=False):
        # Brand statement footer: "for every child, <keyword>" — tagline Regular + keyword Bold
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(6.95), Inches(6), Inches(0.4))
        tf = tb.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        r1 = p.add_run(); r1.text = "for every child, "
        r1.font.name = MINOR_FONT; r1.font.size = Pt(12); r1.font.bold = False
        r2 = p.add_run(); r2.text = "every right"
        r2.font.name = MINOR_FONT; r2.font.size = Pt(12); r2.font.bold = True
        col = white if dark_bg else blue
        for r in (r1, r2):
            r.font.color.rgb = col

    def paint_bg(slide, hexval):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor.from_string(hexval)

    # --- Slide 1: Title (UNICEF Blue full-bleed cover) ---
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    paint_bg(s, UNICEF_BLUE)
    tb = s.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(11.7), Inches(2.5))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = "Presentation title"
    p.font.name = MAJOR_FONT; p.font.size = Pt(44); p.font.bold = True
    p.font.color.rgb = white
    p2 = tf.add_paragraph(); p2.text = "Subtitle / author / date"
    p2.font.name = MINOR_FONT; p2.font.size = Pt(20); p2.font.color.rgb = white
    add_footer(s, dark_bg=True)

    # --- Slide 2: Section divider ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint_bg(s, WHITE)
    bar = s.shapes.add_shape(1, Inches(0), Inches(3.0), Inches(0.35), Inches(1.6))
    bar.fill.solid(); bar.fill.fore_color.rgb = blue; bar.line.fill.background()
    tb = s.shapes.add_textbox(Inches(0.8), Inches(3.1), Inches(11), Inches(1.4))
    p = tb.text_frame.paragraphs[0]; p.text = "Section title"
    p.font.name = MAJOR_FONT; p.font.size = Pt(36); p.font.bold = True
    p.font.color.rgb = blue
    add_footer(s)

    # --- Slide 3: Content ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    paint_bg(s, WHITE)
    # Heading
    hb = s.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.9))
    hp = hb.text_frame.paragraphs[0]; hp.text = "Content heading"
    hp.font.name = MAJOR_FONT; hp.font.size = Pt(28); hp.font.bold = True
    hp.font.color.rgb = blue
    # Underline accent rule
    rule = s.shapes.add_shape(1, Inches(0.8), Inches(1.35), Inches(2.2), Inches(0.06))
    rule.fill.solid(); rule.fill.fore_color.rgb = blue; rule.line.fill.background()
    # Body bullets
    body = s.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(11.7), Inches(4.8))
    bf = body.text_frame; bf.word_wrap = True
    samples = [
        "Body text uses Noto Sans Regular for maximum legibility.",
        "UNICEF Blue is the dominant colour; accents are used sparingly.",
        "Keep layouts clean — avoid visual clutter (Brand Book 4.0, p.4).",
    ]
    for i, txt in enumerate(samples):
        para = bf.paragraphs[0] if i == 0 else bf.add_paragraph()
        para.text = "•  " + txt
        para.font.name = MINOR_FONT; para.font.size = Pt(18)
        para.font.color.rgb = black
        para.space_after = Pt(10)
    add_footer(s)

    tmp_pptx = os.path.join(OUT, "_tmp_deck.pptx")
    prs.save(tmp_pptx)
    inject_theme(tmp_pptx)
    to_template(tmp_pptx, path, "pptx")
    os.remove(tmp_pptx)
    return path


# ---------------------------------------------------------------------------
# 3. Word template (.dotx)
# ---------------------------------------------------------------------------
def _set_style_font(style, name, size_pt, bold, color_hex, space_after=8):
    f = style.font
    f.name = name
    f.size = DPt(size_pt)
    f.bold = bold
    f.color.rgb = DRGB.from_string(color_hex)
    # ensure east-asian/complex script also map to the font
    rpr = style.element.get_or_add_rPr()
    rfonts = rpr.find(qn("w:rFonts"))
    if rfonts is None:
        rfonts = rpr.makeelement(qn("w:rFonts"), {})
        rpr.append(rfonts)
    for attr in ("w:ascii", "w:hAnsi", "w:cs"):
        rfonts.set(qn(attr), name)
    pf = style.paragraph_format
    pf.space_after = DPt(space_after)


def build_dotx(path):
    doc = Document()

    styles = doc.styles
    # Normal / body
    _set_style_font(styles["Normal"], MINOR_FONT, 11, False, BLACK, space_after=8)
    styles["Normal"].paragraph_format.line_spacing = 1.15

    # Headings
    _set_style_font(styles["Title"], MAJOR_FONT, 32, True, UNICEF_BLUE, space_after=6)
    _set_style_font(styles["Heading 1"], MAJOR_FONT, 20, True, UNICEF_BLUE, space_after=6)
    _set_style_font(styles["Heading 2"], MAJOR_FONT, 15, True, ACC_BLUE, space_after=4)
    _set_style_font(styles["Heading 3"], MAJOR_FONT, 12, True, ACC_GREEN, space_after=4)

    # Caption + Quote if present
    for name, size, color in (("Caption", 9, ACC_GREEN),):
        if name in [s.name for s in styles]:
            _set_style_font(styles[name], MINOR_FONT, size, False, color)

    # Sample content so the user sees the system applied
    doc.add_paragraph("Document title", style="Title")
    p = doc.add_paragraph("for every child, ")
    p.runs[0].font.name = MINOR_FONT
    r = p.add_run("every right"); r.bold = True; r.font.name = MINOR_FONT
    for run in p.runs:
        run.font.color.rgb = DRGB.from_string(UNICEF_BLUE)
        run.font.size = DPt(12)

    doc.add_paragraph("Heading 1", style="Heading 1")
    doc.add_paragraph(
        "Body text is set in Noto Sans Regular at 11pt. UNICEF Blue is reserved "
        "for headings and emphasis and remains the dominant brand colour. Accent "
        "colours are used sparingly and never take prominence over UNICEF Blue "
        "(Brand Book 4.0, p.53).")
    doc.add_paragraph("Heading 2", style="Heading 2")
    doc.add_paragraph(
        "Colours may not be lightened, darkened, or shown transparently; tints are "
        "permitted only in data visualization and publication inside-pages (p.53).")
    doc.add_paragraph("Heading 3", style="Heading 3")
    doc.add_paragraph("Body paragraph under a level-3 heading.")

    tmp_docx = os.path.join(OUT, "_tmp_doc.docx")
    doc.save(tmp_docx)
    inject_theme(tmp_docx)
    to_template(tmp_docx, path, "docx")
    os.remove(tmp_docx)
    return path


if __name__ == "__main__":
    t = build_thmx(os.path.join(OUT, "UNICEF-Brand-4.0.thmx"))
    p = build_potx(os.path.join(OUT, "UNICEF-Brand-4.0.potx"))
    d = build_dotx(os.path.join(OUT, "UNICEF-Brand-4.0.dotx"))
    for f in (t, p, d):
        print(f"  wrote {os.path.relpath(f, HERE)}  ({os.path.getsize(f):,} bytes)")
    print("Done.")
